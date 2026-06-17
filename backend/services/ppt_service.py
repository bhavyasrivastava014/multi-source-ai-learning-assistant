from pptx import Presentation


def extract_text_from_ppt(file_path: str) -> str:
    presentation = Presentation(file_path)
    text = ""

    for slide_num, slide in enumerate(presentation.slides, start=1):
        text += f"\n--- Slide {slide_num} ---\n"

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text